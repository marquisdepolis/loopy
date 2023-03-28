# %%
import openai
import os
from time import time, sleep
import numpy as np
import json
import textwrap
import pandas as pd
import re
from sys import exit
import PyPDF2
# import PptxReader
import pptx
# import tkinter as tk
# from tkinter import filedialog


def open_file(filepath):
    with open(filepath, 'r', encoding='latin-1') as infile:
        return infile.read()


os.environ["OPENAI_API_KEY"] = open_file('openai_api_key.txt')
openai.api_key = open_file('openai_api_key.txt')
openai_api_key = openai.api_key

# %%


def inputfile():
    # input file
    fp = filedialog.askopenfilename()
    ext = os.path.splitext(fp)[-1].lower()

    # use == to check for equality
    if ext == ".pptx":
        print(fp, "is a pptx!")
        PptxReader = download_loader("PptxReader")
        loader = PptxReader()
        documents = loader.load_data(file=Path(fp))
        ppt = pptx.Presentation(fp)
        # Loop through each slide and extract the text
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    print(shape.text)
    elif ext == ".pdf":
        print(fp, "is a pdf file!")
        reader = PyPDF2.PdfReader(fp)
        alltext = reader.pages[0].extract_text()
    elif ext == ".html":
        with open(fp, 'r') as f:
            reader = BeautifulSoup(f, 'html.parser')
            alltext = reader.get_text()
    elif ext == ".epub":
        print(fp, "is an epub file!")
        alltext = extract_epub_contents(fp)
    else:
        print(fp, "is an unknown file format.")


def truncate_text(text, max_tokens):
    tokens = text.split(" ")
    if len(tokens) > max_tokens:
        tokens = tokens[:max_tokens]
    return " ".join(tokens)

# %%
# build the index


def gpt3_embedding(content, engine='text-embedding-ada-002'):
    try:
        content = content.encode(encoding='ASCII', errors='ignore').decode()
    except UnicodeEncodeError:
        content = content.encode(encoding='utf-8', errors='ignore').decode()
    response = openai.Embedding.create(input=content, engine=engine)
    vector = response['data'][0]['embedding']
    return vector


# Only for writing into the index
# if __name__ == '__main__':
#     chunks = textwrap.wrap(alltext, 3000)
#     result = list()
#     count = 0
#     for chunk in chunks:
#         count = count + 1
#         embedding = gpt3_embedding(chunk.encode(
#             encoding='ASCII', errors='ignore').decode())
#         info = {'content': chunk, 'vector': embedding}
#         print(info, '\n\n\n')
#         result.append(info)
#     with open('tsla_index.json', 'w') as outfile:
#         json.dump(result, outfile, indent=2)


# %%
# return dot product of two vectors
def similarity(v1, v2):
    return np.dot(v1, v2)


def search_index(text, data, count=10):
    embed = gpt3_embedding(truncate_text(text, 6000))
    scores = list()
    for item in data:
        # Only needed because of token considerations
        temptext = truncate_text(item['text'], 6000)
        score = similarity(embed, gpt3_embedding(temptext))
        scores.append({'content': temptext, 'score': score})
    ordered = sorted(scores, key=lambda d: d['score'], reverse=True)
    return ordered[0:count]


def gpt3_completion(prompt, engine='text-davinci-003', temp=0.6, top_p=1.0, tokens=2000, freq_pen=0.25, pres_pen=0.0, stop=['<<END>>']):
    max_retry = 5
    retry = 0
    prompt = prompt.encode(encoding='ASCII', errors='ignore').decode()
    while True:
        try:
            response = openai.Completion.create(
                engine=engine,
                prompt=prompt,
                temperature=temp,
                max_tokens=tokens,
                top_p=top_p,
                frequency_penalty=freq_pen,
                presence_penalty=pres_pen,
                stop=stop)
            text = response['choices'][0]['text'].strip()
            text = re.sub('\s+', ' ', text)
            filename = '%s_gpt3.txt' % time()
            with open('gpt3_logs/%s' % filename, 'w') as outfile:
                outfile.write('PROMPT:\n\n' + prompt +
                              '\n\n==========\n\nRESPONSE:\n\n' + text)
            return text
        except Exception as oops:
            retry += 1
            if retry >= max_retry:
                return "GPT3 error: %s" % oops
            print('Error communicating with OpenAI:', oops)
            sleep(1)


def inputquery(query):
    with open('slcindex.json', 'r') as infile:
        data = json.load(infile)
    data_list = list(data['docstore']['docs'].values())
    results = search_index(query, data_list)
    answers = list()
    for result in results:
        answer = gpt3_completion(query)
        print('\n\n', answer)
        answers.append(answer)
    # summarise answers together
    all_answers = '\n\n'.join(answers)
    chunks = textwrap.wrap(all_answers, 10000)
    final = list()
    for chunk in chunks:
        prompt = "Please summarise:" + chunk
        summary = gpt3_completion(prompt)
        final.append(summary)
    return final


# inputquery("what did rohit write about strange equation in the canon")
