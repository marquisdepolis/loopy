# %%
import pptx
import os
from llama_index import download_loader
from pathlib import Path
import PyPDF2
from pptx import Presentation


def read_file(fp):
    ext = os.path.splitext(fp)[-1].lower()
    if ext == ".pptx":
        print(fp, "is a pptx!")
        presentation = Presentation(fp)
        textshape = []
        for slide in presentation.slides:
            # Iterate through the shapes (text boxes, tables, images, etc.) in the slide
            for shape in slide.shapes:
                # Check if the shape has a text frame
                if hasattr(shape, "text"):
                    # Append the text from the text frame to the texts list
                    textshape.append(shape.text)

        # Join the texts list into a single string, separated by line breaks
        texts = "\n".join(textshape)
    elif ext == ".pdf":
        print(fp, "is a pdf file!")
        reader = PyPDF2.PdfReader(fp)
        texts = []
        # Loop through each page and extract the text
        number_of_pages = len(reader.pages)
        for x in range(number_of_pages):
            page = reader.pages[x]
            text = page.extract_text()
            texts.append(text)
    elif ext == ".html":
        print(fp, "is an html file!")
        with open(fp, 'r') as f:
            reader = BeautifulSoup(f, 'html.parser')
            texts = reader.get_text()
    elif filename.endswith('.epub'):
        texts = []
        book = epub.read_epub(fp)
        chapters = book.get_items_of_type(ebooklib.ITEM_DOCUMENT)
        for chapter in chapters:
            soup = BeautifulSoup(chapter.get_content(), 'html.parser')
            text = soup.get_text()
            texts.append(text)
    else:
        print(fp, "is an unknown file format.")

    return texts
